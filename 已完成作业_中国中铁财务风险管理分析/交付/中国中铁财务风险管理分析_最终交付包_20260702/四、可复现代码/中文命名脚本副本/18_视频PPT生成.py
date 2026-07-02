"""Build a 3-minute presentation deck for the course video."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "docs/assets/figures"
OUTPUT = ROOT / "outputs/video/china_railway_risk_3min_deck.pptx"

TITLE = "中国中铁财务风险管理分析"
SUBTITLE = "风险图谱、文本指标与机器学习预警"

SLIDES = [
    {
        "title": TITLE,
        "subtitle": SUBTITLE,
        "bullets": ["公开数据", "风险传导", "预警模型", "弹性管理"],
    },
    {
        "title": "经营特征：项目链条拉长资金压力",
        "image": "financial_trends.png",
        "bullets": ["收入和利润在 2023 年后承压", "合同资产与应收账款占用上升", "资产负债率维持高位"],
    },
    {
        "title": "数据口径：把财务、文本和事件放进同一框架",
        "bullets": ["年报、公告和评级报告", "裁判文书、执行信息和企业风险线索", "同业建筑企业财务面板"],
        "flow": ["公开披露", "风险事件", "风险图谱", "文本指标", "预警模型"],
    },
    {
        "title": "风险图谱：识别高中心性节点",
        "image": "risk_network_gephi_style.png",
        "bullets": ["77 个节点、133 条边", "合规风险与流动性风险处于核心位置", "诉讼和执行事件具有桥接作用"],
    },
    {
        "title": "文本与事件：软信号补充财务指标",
        "images": ["text_risk_heatmap.png", "risk_event_matrix.png"],
        "bullets": ["偿债、组织传导、市场和项目风险长期高关注", "合同资产占用和偿债压力位于高风险区"],
    },
    {
        "title": "机器学习预警：把同业财务面板转为压力信号",
        "bullets": ["样本：11 家建筑企业、2021-2025 年财务面板", "模型：Logistic Regression 与 Random Forest", "中国中铁 2026 年压力观察概率偏高"],
        "flow": ["财务面板", "规则标签", "模型训练", "概率输出", "管理解释"],
    },
    {
        "title": "弹性管理：判断缓冲能力的短板",
        "images": ["resilience_radar_2025.png", "resilience_score_trend.png"],
        "bullets": ["2025 年治理信用缓冲仍有支撑", "经营缓冲和财务缓冲偏弱", "重点投向回款、合同资产和高中心性节点治理"],
    },
    {
        "title": "结论：形成识别、评估和预警闭环",
        "bullets": ["穿透监测合同资产和应收账款", "跟踪债务覆盖与融资期限结构", "复核子公司诉讼、执行和合规事件", "用图谱中心性确定治理优先级"],
    },
]


def add_textbox(slide, left, top, width, height, text, size=24, bold=False, color=RGBColor(30, 35, 40)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, left=Inches(0.7), top=Inches(5.65), width=Inches(12.0), height=Inches(1.0)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = RGBColor(55, 61, 70)
    return box


def add_flow(slide, items):
    left = Inches(0.75)
    top = Inches(2.7)
    box_w = Inches(2.1)
    gap = Inches(0.25)
    for idx, item in enumerate(items):
        shape = slide.shapes.add_shape(1, left + idx * (box_w + gap), top, box_w, Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(234, 242, 250)
        shape.line.color.rgb = RGBColor(42, 102, 143)
        frame = shape.text_frame
        frame.text = item
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(25, 65, 95)
        if idx < len(items) - 1:
            add_textbox(slide, left + (idx + 1) * box_w + idx * gap, top + Inches(0.18), gap, Inches(0.4), "→", size=20, bold=True)


def add_image(slide, name, left=Inches(0.75), top=Inches(1.35), width=Inches(11.8), height=Inches(4.1)):
    path = FIGURES / name
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        add_textbox(slide, Inches(0.65), Inches(0.35), Inches(12), Inches(0.55), spec["title"], size=24, bold=True)
        if idx == 0:
            add_textbox(slide, Inches(0.8), Inches(1.45), Inches(12), Inches(0.45), spec["subtitle"], size=20)
            add_flow(slide, spec["bullets"])
            add_textbox(slide, Inches(0.8), Inches(5.55), Inches(11.6), Inches(0.5), "3 分钟展示：风险传导 -> 风险评估 -> 风险预警", size=18, bold=True)
            continue
        if "image" in spec:
            add_image(slide, spec["image"])
        if "images" in spec:
            add_image(slide, spec["images"][0], left=Inches(0.65), top=Inches(1.25), width=Inches(6.1), height=Inches(4.0))
            add_image(slide, spec["images"][1], left=Inches(6.9), top=Inches(1.25), width=Inches(5.75), height=Inches(4.0))
        if "flow" in spec:
            add_flow(slide, spec["flow"])
        add_bullets(slide, spec["bullets"])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"wrote {OUTPUT}")


if __name__ == "__main__":
    build_deck()
